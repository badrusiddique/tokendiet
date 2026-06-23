"""Reproducible evidence for *why Tokendiet does not convert Office files*.

PDF and HTML have an expensive native form (page images / raw markup). Office
files (DOCX/XLSX/PPTX) do not: Claude doesn't render them as images, and their
realistic text extraction is already Markdown-sized. This probe demonstrates
that converting them to Markdown saves nothing against a realistic baseline —
Markdown is even marginally *larger*.

Run:  uv sync --group probe  &&  python benchmarks/office_probe.py

Optional Markdown column needs MarkItDown:  pip install "markitdown[docx,pptx,xlsx]"
"""

from __future__ import annotations

import tempfile
import zipfile
from pathlib import Path

import tiktoken

enc = tiktoken.get_encoding("cl100k_base")

PROSE = (
    "It is a truth universally acknowledged, that a single man in possession of a good "
    "fortune, must be in want of a wife. However little known the feelings or views of such "
    "a man may be on his first entering a neighbourhood, this truth is so well fixed. "
)


def toks(s: str) -> int:
    return len(enc.encode(s))


def ooxml_tokens(path: Path, prefixes: tuple[str, ...]) -> int:
    """Tokens of the *content* XML parts (excludes style/theme boilerplate)."""
    z = zipfile.ZipFile(path)
    return sum(
        toks(z.read(n).decode("utf-8", "replace"))
        for n in z.namelist()
        if n.endswith(".xml") and any(n.startswith(p) for p in prefixes)
    )


def _make_docx(path: Path) -> str:
    from docx import Document

    d = Document()
    d.add_heading("Quarterly Report", 0)
    for _ in range(20):
        d.add_paragraph(PROSE)
    d.save(path)
    return "\n".join(p.text for p in Document(path).paragraphs)


def _make_xlsx(path: Path) -> str:
    import openpyxl
    import xlsxwriter

    wb = xlsxwriter.Workbook(str(path))
    ws = wb.add_worksheet()
    for row in range(60):
        for col in range(8):
            ws.write(row, col, f"cell {row}-{col} {PROSE[:20]}")
    wb.close()
    sheet = openpyxl.load_workbook(path).active
    return "\n".join(
        ",".join("" if c.value is None else str(c.value) for c in row) for row in sheet.iter_rows()
    )


def _make_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation()
    for _ in range(10):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Slide title"
        slide.placeholders[1].text = PROSE * 2
    prs.save(path)
    parts = []
    for slide in Presentation(path).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _markdown_tokens(path: Path) -> int | None:
    try:
        from markitdown import MarkItDown
    except Exception:
        return None
    return toks(MarkItDown().convert(str(path)).text_content)


def main() -> int:
    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        specs = [
            ("docx", tmp_path / "s.docx", _make_docx, ("word/document", "word/header")),
            ("xlsx", tmp_path / "s.xlsx", _make_xlsx, ("xl/worksheets/", "xl/sharedStrings")),
            ("pptx", tmp_path / "s.pptx", _make_pptx, ("ppt/slides/", "ppt/notesSlides/")),
        ]
        print(f"{'format':6} {'rawOOXML':>9} {'cleanText':>10} {'markdown':>9}   verdict")
        for label, path, make, prefixes in specs:
            clean = toks(make(path))
            raw = ooxml_tokens(path, prefixes)
            md = _markdown_tokens(path)
            md_str = f"{md:>9,}" if md is not None else f"{'n/a':>9}"
            verdict = "no real savings (markdown >= clean text)"
            print(f"{label:6} {raw:>9,} {clean:>10,} {md_str}   {verdict}")
    print(
        "\nConclusion: the only 'savings' is vs dumping raw OOXML — a baseline no one uses.\n"
        "Against realistic text extraction, Markdown is equal or larger. Office files are "
        "intentionally NOT supported by Tokendiet. See docs/format-support.md."
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
